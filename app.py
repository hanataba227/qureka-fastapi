from openai import OpenAI
import fitz  # PDF
import os
from dotenv import load_dotenv
from pptx import Presentation  # PPTX

# 환경변수 로딩
load_dotenv()
api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)

# PDF 텍스트 추출
def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

# PPT 텍스트 추출
def extract_text_from_pptx(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# 요약 프롬프트 설정
def get_summary_prompt(summary_type, content):
    prompts = {
        "기본 요약": {
            "system": f'''내용''',
            "user": f'''내용:\n\n{content}'''
        },
        "핵심 요약": {
            "system": f'''내용''',
            "user": f'''내용:\n\n{content}'''
        },
        "주제 요약": {
            "system": f'''내용''',
            "user": f'''내용:\n\n{content}'''
        },
        "목차 요약": {
            "system": f'''내용''',
            "user": f'''내용:\n\n{content}'''
        },
        "키워드 요약": {
            "system": f'''내용''',
            "user": f'''내용:\n\n{content}'''
        }
    }
    return prompts.get(summary_type, {"system": "", "user": ""})


# GPT API 호출
def summarize_with_chatgpt(system_message, user_message):
    completion = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_message},
            {"role": "user", "content": user_message}
        ],
        temperature=0.7
    )
    summary = completion.choices[0].message.content
    usage = completion.usage
    return summary, usage

# 실행 함수
def run_summary(file_path, summary_type):
    print(f">>> [{summary_type}] 요약 시작...")

    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        content = extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        content = extract_text_from_pptx(file_path)
    else:
        print("지원되지 않는 파일 형식입니다. PDF 또는 PPTX만 가능합니다.")
        return

    prompt = get_summary_prompt(summary_type, content)
    if not prompt or not prompt["system"] or not prompt["user"]:
        print("요약 유형이 잘못되었거나 프롬프트 구성에 문제가 있습니다.")
        return

    summary, usage = summarize_with_chatgpt(prompt["system"], prompt["user"])

    print("\n요약 결과:\n")
    print(summary)
    print("\n사용한 토큰 수:")
    print(f"- prompt_tokens: {usage.prompt_tokens}")
    print(f"- completion_tokens: {usage.completion_tokens}")
    print(f"- total_tokens: {usage.total_tokens}")

# 실행 예시
if __name__ == "__main__":
    file_path = ".pdf"  # 분석 대상 경로
    summary_type = "기본 요약"  # 기본/핵심/주제/목차/키워드 중 택1
    run_summary(file_path, summary_type)
